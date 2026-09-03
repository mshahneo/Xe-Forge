"""Build docs/mlir_backend_overview.pptx from the 5-slide overview.

Run:  uv run --with python-pptx python docs/build_overview_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (Intel-ish blue, dark slate text, neutral accents) ----
INK      = RGBColor(0x1A, 0x1F, 0x2B)   # near-black slate
BLUE     = RGBColor(0x00, 0x68, 0xB5)   # intel blue
BLUE_DK  = RGBColor(0x00, 0x3C, 0x71)
ACCENT   = RGBColor(0x00, 0xA3, 0xA1)   # teal
MUTED    = RGBColor(0x5B, 0x66, 0x77)
LIGHT    = RGBColor(0xF3, 0xF6, 0xFA)
CODE_BG  = RGBColor(0x11, 0x18, 0x22)
CODE_FG  = RGBColor(0xE6, 0xED, 0xF3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RULE     = RGBColor(0xD8, 0xDF, 0xE8)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

MONO = "Consolas"
SANS = "Calibri"


def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def run(p, text, size, color=INK, bold=False, italic=False, font=SANS):
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return r


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def header(slide, kicker, title):
    """Top band with kicker + title, and a rule under it."""
    rect(slide, 0, 0, SW, Inches(1.45), LIGHT)
    rect(slide, 0, Inches(1.45), SW, Pt(3), BLUE)
    # accent tab
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.16), Inches(0.62), ACCENT)
    tf = textbox(slide, Inches(0.85), Inches(0.30), Inches(11.9), Inches(1.05))
    p = para(tf, first=True)
    run(p, kicker.upper(), 12, BLUE, bold=True)
    p.space_after = Pt(2)
    p2 = tf.add_paragraph()
    run(p2, title, 26, INK, bold=True)


def body_box(slide, top=Inches(1.75), height=Inches(5.35)):
    return textbox(slide, Inches(0.85), top, Inches(11.65), height)


def bullet(tf, runs, level=0, size=15, space_after=8, first=False, bullet_char="▸"):
    """runs: list of (text, kwargs) tuples."""
    p = para(tf, first=first)
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if bullet_char:
        run(p, bullet_char + "  ", size, BLUE if level == 0 else ACCENT, bold=True)
    for text, kw in runs:
        run(p, text, kw.get("size", size), kw.get("color", INK),
            kw.get("bold", False), kw.get("italic", False), kw.get("font", SANS))
    return p


# =====================================================================
# Slide 0 — Title
# =====================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.35), SH, BLUE)
rect(s, Inches(0.35), 0, Inches(0.10), SH, ACCENT)
tf = textbox(s, Inches(1.1), Inches(2.1), Inches(11.3), Inches(3.3), anchor=MSO_ANCHOR.TOP)
p = para(tf, first=True)
run(p, "MLIR / XeGPU BACKEND FOR XE-FORGE", 14, ACCENT, bold=True)
p.space_after = Pt(10)
p2 = tf.add_paragraph()
run(p2, "Bringing the Xe-Forge optimization loop", 40, WHITE, bold=True)
p3 = tf.add_paragraph()
run(p3, "to MLIR on Intel Xe2 “Battlemage”", 40, WHITE, bold=True)
p3.space_after = Pt(18)
p4 = tf.add_paragraph()
run(p4, "A 5-part overview — same LLM-driven loop, a new DSL plugged in", 18, RGBColor(0xB6,0xC2,0xD2))
# footer
tf2 = textbox(s, Inches(1.1), Inches(6.5), Inches(11.3), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
pf = para(tf2, first=True)
run(pf, "Companion to docs/mlir_backend_architecture.md", 12, MUTED, italic=True)

# =====================================================================
# Slide 1 — The Idea
# =====================================================================
s = add_slide()
header(s, "Slide 1 · The Idea", "A DSL-agnostic loop, one new engine")
tf = body_box(s)
bullet(tf, [("Xe-Forge, originally:  ", {"bold": True}),
            ("an LLM-driven optimizer for ", {}),
            ("Triton", {"bold": True, "color": BLUE}),
            (" kernels on Intel XPU — a multi-stage loop:", {})],
       first=True, space_after=4)
p = tf.add_paragraph(); p.level = 1; p.space_after = Pt(10)
run(p, "analyze → plan → optimize (verify & retry) → gate on correctness + speedup", 15, MUTED, italic=True, font=MONO)
bullet(tf, [("Driven by ", {}), ("CoVeR", {"bold": True, "color": BLUE}),
            (" (Chain of Verification & Refinement) + an Intel-GPU ", {}),
            ("Knowledge Base", {"bold": True, "color": BLUE}), (".", {})], space_after=14)

bullet(tf, [("What we added:  ", {"bold": True}),
            ("a fifth DSL — ", {}), ("MLIR / XeGPU", {"bold": True, "color": ACCENT}),
            (", targeting Intel Xe2 “Battlemage”.", {})], space_after=14)

bullet(tf, [("The bet that paid off:  ", {"bold": True}),
            ("the optimization loop is ", {}), ("DSL-agnostic", {"bold": True}),
            (". MLIR is a ", {}), ("plug-in, not a fork", {"bold": True, "italic": True}),
            (" — it reuses the analyzer, planner, CoVeR, gating, trial tree and KB "
             "framework unchanged, and only swaps the executor and a few DSPy signatures.", {})],
       space_after=16)

# pull-quote band
qb = rect(s, Inches(0.85), Inches(5.95), Inches(11.65), Inches(0.95), LIGHT)
rect(s, Inches(0.85), Inches(5.95), Pt(5), Inches(0.95), ACCENT)
tq = textbox(s, Inches(1.15), Inches(5.95), Inches(11.1), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
pq = para(tq, first=True)
run(pq, "“Same driver and race strategy; we swapped the engine and fuel.”", 17, BLUE_DK, italic=True, bold=True)

# =====================================================================
# Slide 2 — What's the Same (table)
# =====================================================================
s = add_slide()
header(s, "Slide 2 · What's the Same", "MLIR inherits the entire Xe-Forge spine")
rows = [
    ("CoVeR verify-and-revise loop", "takes a signature + a verify tool; doesn't care about DSL"),
    ("PlannerAgent (issue → stage plan)", "fully DSL-agnostic, zero DSL references"),
    ("Analyzer / Optimizer agents", "hold 3 signatures each; pick the MLIR one by self.dsl"),
    ("Stage enum, issue taxonomy, best-of-k, gating", "shared code  (_MIN_IMPROVEMENT = 1.02)"),
    ("KB framework (load / format / resolve)", "content scoped per-DSL directory"),
    ("common/ KB (algorithmic + correctness)", "loaded for every DSL, MLIR included"),
]
left = Inches(0.85); top = Inches(1.85)
w_all = Inches(11.65); w1 = Inches(4.9); w2 = w_all - w1
rh = Inches(0.62)
# table header
rect(s, left, top, w1, rh, BLUE)
rect(s, left + w1, top, w2, rh, BLUE_DK)
for x, wdt, txt in [(left, w1, "Reused component"), (left + w1, w2, "Why it just works")]:
    t = textbox(s, x + Inches(0.15), top, wdt - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
    run(para(t, first=True), txt, 15, WHITE, bold=True)
# body rows
for i, (a, b) in enumerate(rows):
    y = top + rh + Emu(int(rh) * i)
    bg = WHITE if i % 2 == 0 else LIGHT
    rect(s, left, y, w1, rh, bg, line=RULE)
    rect(s, left + w1, y, w2, rh, bg, line=RULE)
    ta = textbox(s, left + Inches(0.15), y, w1 - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
    run(para(ta, first=True), a, 13.5, INK, bold=True)
    tb = textbox(s, left + w1 + Inches(0.15), y, w2 - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
    run(para(tb, first=True), b, 13.5, MUTED)

# footnote callout
fy = top + rh + Emu(int(rh) * len(rows)) + Inches(0.12)
cb = rect(s, left, fy, w_all, Inches(0.7), INK)
rect(s, left, fy, Pt(5), Inches(0.7), ACCENT)
tc = textbox(s, left + Inches(0.25), fy, w_all - Inches(0.5), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
pc = para(tc, first=True)
run(pc, "Core cost of adding MLIR:  ", 14, ACCENT, bold=True)
run(pc, "one DSL value (DSL.MLIR) + one new stage (LINALG_LOWERING).", 14, WHITE)

# =====================================================================
# Slide 3 — What's Different (3 cards)
# =====================================================================
s = add_slide()
header(s, "Slide 3 · What's Different", "Only three structural divergences")
cards = [
    ("1", "Self-contained correctness oracle",
     "No PyTorch. The kernel carries its own @main harness with an in-file CPU "
     "reference and prints [ALLCLOSE: TRUE]. Core stays torch-free. The optimizer "
     "edits only the gpu.module / #xegpu.layout / launch geometry — never the harness."),
    ("2", "LINALG_LOWERING — a pre-stage with no Triton analog",
     "MLIR kernels can arrive as high-level linalg and must be lowered to "
     "workgroup-level XeGPU first. Lowering is itself an optimization (tile-search), "
     "so it's a first-class stage that runs before analysis."),
    ("3", "GRF_SWEEP — a correctness-free device knob",
     "Run the same IR with / without the IGC large-register-file flag, keep the "
     "faster. A compile flag, not an IR edit → correctness untouched."),
]
cw = Inches(3.72); gap = Inches(0.24); cx = Inches(0.85); cy = Inches(1.95); ch = Inches(4.35)
for i, (num, title, desc) in enumerate(cards):
    x = cx + Emu(int(cw + gap) * i)
    rect(s, x, cy, cw, ch, WHITE, line=RULE)
    rect(s, x, cy, cw, Inches(0.14), BLUE)          # top accent bar
    # number disc
    from pptx.enum.shapes import MSO_SHAPE
    disc = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.28), cy + Inches(0.42),
                              Inches(0.72), Inches(0.72))
    disc.fill.solid(); disc.fill.fore_color.rgb = ACCENT; disc.line.fill.background()
    disc.shadow.inherit = False
    dtf = disc.text_frame; dtf.word_wrap = False
    dp = dtf.paragraphs[0]; dp.alignment = PP_ALIGN.CENTER
    run(dp, num, 26, WHITE, bold=True)
    # title
    tt = textbox(s, x + Inches(0.28), cy + Inches(1.32), cw - Inches(0.56), Inches(1.2))
    run(para(tt, first=True), title, 16, BLUE_DK, bold=True)
    # desc
    td = textbox(s, x + Inches(0.28), cy + Inches(2.5), cw - Inches(0.56), Inches(1.7))
    run(para(td, first=True), desc, 12.5, MUTED)

tf = textbox(s, Inches(0.85), Inches(6.55), Inches(11.65), Inches(0.5))
p = para(tf, first=True)
run(p, "Everything else in the loop is the same code the Triton path runs.", 15, INK, italic=True, bold=True)

# =====================================================================
# Slide 4 — Two-Level Flow
# =====================================================================
s = add_slide()
header(s, "Slide 4 · The Two-Level MLIR Flow", "Lower first, then optimize")

# flow pipeline
from pptx.enum.shapes import MSO_SHAPE
stages = [("linalg.matmul", "high level", LIGHT, INK),
          ("LEVEL 1: LOWER", "tile-search", BLUE, WHITE),
          ("xegpu WG kernel", "@main + [ALLCLOSE]", LIGHT, INK),
          ("LEVEL 2: OPTIMIZE", "shared loop", ACCENT, WHITE)]
fx = Inches(0.85); fy = Inches(1.85); bw = Inches(2.62); bh = Inches(1.0); fgap = Inches(0.42)
for i, (t1, t2, bg, fg) in enumerate(stages):
    x = fx + Emu(int(bw + fgap) * i)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, fy, bw, bh)
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = RULE; box.line.width = Pt(1); box.shadow.inherit = False
    bt = box.text_frame; bt.word_wrap = True
    bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    run(bp, t1, 14, fg, bold=True, font=MONO)
    bp2 = bt.add_paragraph(); bp2.alignment = PP_ALIGN.CENTER
    run(bp2, t2, 10.5, fg, italic=True)
    if i < len(stages) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw + Emu(int(fgap*0.1)),
                                fy + Inches(0.34), Inches(0.34), Inches(0.32))
        ar.fill.solid(); ar.fill.fore_color.rgb = MUTED; ar.line.fill.background()
        ar.shadow.inherit = False

# Level 1 detail
tf = textbox(s, Inches(0.85), Inches(3.15), Inches(11.65), Inches(2.4))
p = para(tf, first=True); p.space_after = Pt(4)
run(p, "Level 1 — Lowering ", 16, BLUE, bold=True)
run(p, "(_maybe_lower_linalg), routed by structure:", 14, INK)
bullet(tf, [("matmul / MLP", {"bold": True}),
            (" → native 3-stage lowering (tile+vectorize → outline+attach-target → WG "
             "layout), a 5-knob LoweringConfig. Tile-search is ", {}),
            ("hybrid", {"bold": True, "color": ACCENT}),
            (": LLM shortlists KB-seeded configs → each timed on GPU → fastest correct "
             "wins, with a lighthouse second opinion.", {})], level=1, size=13.5, space_after=6)
bullet(tf, [("attention / softmax / layer-norm", {"bold": True}),
            (" → lighthouse dump-only backend (upstream llvm/lighthouse schedules).", {})],
       level=1, size=13.5, space_after=10)
p = tf.add_paragraph(); p.space_after = Pt(4)
run(p, "Level 2 — Optimize:  ", 16, ACCENT, bold=True)
run(p, "the standard analyze → plan → CoVeR-optimize → gate loop runs on the WG kernel.", 14, INK)

# hardening callout
hy = Inches(6.25)
rect(s, Inches(0.85), hy, Inches(11.65), Inches(0.85), INK)
rect(s, Inches(0.85), hy, Pt(5), Inches(0.85), BLUE)
tc = textbox(s, Inches(1.15), hy, Inches(11.1), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
pc = para(tc, first=True)
run(pc, "Correctness hardening:  ", 13.5, BLUE, bold=True, font=SANS)
run(pc, "a lowered-IR-equivalence guard rejects no-op edits — identical lowered IR ⇒ "
        "dead-code noise (caught a phantom 1.68× from a DCE-eliminated prefetch).", 13.5, WHITE)

# =====================================================================
# Slide 5 — KB & Where Next
# =====================================================================
s = add_slide()
header(s, "Slide 5 · Knowledge Base & Where Next", "Layered KB, and the roadmap")

# left column: KB layering
lx = Inches(0.85); ly = Inches(1.9); lw = Inches(6.35)
tl = textbox(s, lx, ly, lw, Inches(4.9))
p = para(tl, first=True); p.space_after = Pt(6)
run(p, "KB layering", 17, BLUE, bold=True)
run(p, "   common/ → <dsl>/common/ → <dsl>/<device>/", 12.5, MUTED, font=MONO)
bullet(tl, [("Shared today:", {"bold": True}),
            (" common/algorithmic_patterns.yaml, correctness.yaml — MLIR gets them free.", {})],
       size=13, space_after=7)
bullet(tl, [("MLIR-specific:", {"bold": True}),
            (" mlir/xpu/* (WG DPAS / layout / memory) and mlir/linalg/* (lowering seeds).", {})],
       size=13, space_after=7)
bullet(tl, [("New KB feature (reusable):", {"bold": True}),
            (" a precondition field so a pattern (e.g. xegpu_add_prefetch_nd) fires only "
             "when it applies.", {})], size=13, space_after=7)
bullet(tl, [("Overlap", {"bold": True}),
            (" (SIMD16, f32-accum, DPAS shapes, GRF) is re-expressed per DSL — a refactor "
             "opportunity.", {})], size=13, space_after=7)

# right column: forward paths (numbered)
rx = Inches(7.45); rw = Inches(5.05)
rect(s, rx, ly, rw, Inches(4.05), LIGHT)
rect(s, rx, ly, rw, Inches(0.5), BLUE_DK)
th = textbox(s, rx + Inches(0.25), ly, rw - Inches(0.5), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
run(para(th, first=True), "Forward paths (by leverage)", 15, WHITE, bold=True)
paths = [
    "Real xegpu.prefetch_nd for vector-dialect kernels (make transfer_reads lower to load_nd).",
    "Refactor shared hardware facts into common/hardware_xe.yaml.",
    "Native lowering for attention / softmax / layer-norm (off lighthouse).",
    "MLIR support in the ReAct strategy (currently CoVeR-only).",
]
tp = textbox(s, rx + Inches(0.28), ly + Inches(0.7), rw - Inches(0.56), Inches(3.2))
for i, txt in enumerate(paths):
    p = para(tp, first=(i == 0)); p.space_after = Pt(9)
    run(p, f"{i+1}.  ", 13.5, ACCENT, bold=True)
    run(p, txt, 13, INK)

# bottom-line band
by = Inches(6.15)
rect(s, Inches(0.85), by, Inches(11.65), Inches(1.0), INK)
rect(s, Inches(0.85), by, Pt(5), Inches(1.0), ACCENT)
tb = textbox(s, Inches(1.15), by, Inches(11.1), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
pb = para(tb, first=True)
run(pb, "Bottom line:  ", 15, ACCENT, bold=True)
run(pb, "the Xe-Forge loop generalizes beyond Triton. MLIR runs end-to-end on real GPU "
        "— two-level flow, self-contained oracle, verified speedups — for one DSL value, "
        "one new stage, a dedicated executor, and DSL-scoped KB.", 14.5, WHITE)

out = "docs/mlir_backend_overview.pptx"
prs.save(out)
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
